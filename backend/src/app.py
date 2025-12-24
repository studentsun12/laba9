from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFilter
from io import BytesIO
import os
import traceback
from datetime import datetime
import time
import random
import math

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

THEMES = {
    "modern": {
        "primary": (26, 32, 46),
        "secondary": (59, 130, 246),
        "accent": (99, 102, 241),
        "text_main": (255, 255, 255),
        "text_secondary": (229, 231, 235),
    },
    "professional": {
        "primary": (15, 23, 42),
        "secondary": (51, 65, 85),
        "accent": (3, 102, 214),
        "text_main": (255, 255, 255),
        "text_secondary": (203, 213, 225),
    },
    "vibrant": {
        "primary": (55, 35, 30),
        "secondary": (236, 72, 153),
        "accent": (168, 85, 247),
        "text_main": (255, 240, 235),
        "text_secondary": (255, 200, 150),
    },
    "ocean": {
        "primary": (15, 32, 48),
        "secondary": (34, 197, 232),
        "accent": (6, 182, 212),
        "text_main": (226, 232, 240),
        "text_secondary": (148, 163, 184),
    },
    "sunset": {
        "primary": (55, 35, 30),
        "secondary": (255, 120, 80),
        "accent": (255, 160, 100),
        "text_main": (255, 240, 235),
        "text_secondary": (255, 200, 150),
    },
}

def generate_professional_image(width=640, height=480, theme_colors=None, seed=None, style="abstract"):
    """Генерирует профессиональные картинки для презентаций"""
    if seed:
        random.seed(seed)

    if theme_colors is None:
        theme_colors = [(59, 130, 246), (99, 102, 241), (30, 41, 59)]

    img = Image.new('RGB', (width, height), color=theme_colors[0])
    draw = ImageDraw.Draw(img, 'RGBA')

    # 1. Красивый градиент фон
    for y in range(height):
        ratio = y / height
        r = int(theme_colors[0][0] * (1 - ratio) + theme_colors[1][0] * ratio)
        g = int(theme_colors[0][1] * (1 - ratio) + theme_colors[1][1] * ratio)
        b = int(theme_colors[0][2] * (1 - ratio) + theme_colors[1][2] * ratio)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    if style == "abstract":
        for i in range(5):
            x = random.randint(-100, width + 100)
            y = random.randint(-100, height + 100)
            size = random.randint(80, 200)
            color = (*theme_colors[2], random.randint(40, 100))
            draw.ellipse([x-size, y-size, x+size, y+size], fill=color, outline=None)

    elif style == "tech":
        grid_color = (*theme_colors[2], 60)
        for i in range(-height, width, 60):
            draw.line([(i, 0), (i + height, height)], fill=grid_color, width=2)
            draw.line([(i, height), (i + height, 0)], fill=grid_color, width=2)

        for x in range(0, width, 120):
            for y in range(0, height, 120):
                draw.ellipse([x-8, y-8, x+8, y+8], fill=theme_colors[2], outline=None)

    elif style == "organic":
        for wave in range(3):
            points = []
            for x in range(0, width + 20, 20):
                y = int(height // 2 + wave * 60 + 30 * math.sin((x + wave * 100) / 80))
                points.append((x, y))

            if len(points) > 1:
                color = (*theme_colors[2], 80)
                draw.polygon(points + [(width, height), (0, height)], fill=color, outline=None)

    elif style == "cards":
        card_color = (*theme_colors[1], 120)
        for i in range(3):
            x = 80 + i * 170
            y = 100 + (i % 2) * 100
            draw.rectangle([x, y, x+150, y+150], fill=card_color, outline=theme_colors[2], width=2)
            draw.ellipse([x+60, y+60, x+90, y+90], fill=theme_colors[2], outline=None)

    # 2. Декоративные линии
    draw.rectangle([0, 0, 8, height], fill=theme_colors[1], outline=None)
    draw.rectangle([0, 0, width, 4], fill=theme_colors[1], outline=None)

    # 3. Диагональные полосы
    stripe_color = (*theme_colors[2], 70)
    for i in range(-height, width, 100):
        draw.line([(i, 0), (i + height, height)], fill=stripe_color, width=3)

    # 4. Угловые элементы
    corner_size = 60
    corner_color = (*theme_colors[2], 100)
    draw.ellipse([width-corner_size*2, -corner_size, width+corner_size, corner_size*2], 
                 fill=corner_color, outline=None)
    draw.ellipse([-corner_size, height-corner_size*2, corner_size*2, height+corner_size], 
                 fill=corner_color, outline=None)

    img = img.filter(ImageFilter.GaussianBlur(radius=0.5))
    return img

def create_title_slide(prs, topic: str, theme: dict):
    """Титульный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme["primary"])

    top_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    top_shape.fill.solid()
    top_shape.fill.fore_color.rgb = RGBColor(*theme["secondary"])
    top_shape.line.color.rgb = RGBColor(*theme["secondary"])

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.text = topic
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*theme["text_main"])
    p.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
    text_frame = subtitle_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Professional Presentation"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(*theme["secondary"])
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    text_frame = date_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = datetime.now().strftime("%B %d, %Y")
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(*theme["text_secondary"])
    p.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title: str, content: str, seed: int, theme: dict, style: str = "abstract"):
    """Слайд с контентом и красивой локальной картинкой"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme["primary"])

    line_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.7))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(*theme["secondary"])
    line_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(4.5), Inches(0.9))
    text_frame = title_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*theme["text_main"])

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.8), Inches(5.4))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = content
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(*theme["text_secondary"])
    p.line_spacing = 1.6

    try:
        theme_colors = [theme["primary"], theme["secondary"], theme["accent"]]
        img = generate_professional_image(640, 480, theme_colors, seed=seed, style=style)

        img_byte_arr = BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)

        slide.shapes.add_picture(img_byte_arr, Inches(5.3), Inches(0.8), width=Inches(4.2), height=Inches(3.15))

        frame = slide.shapes.add_shape(1, Inches(5.3), Inches(0.8), Inches(4.2), Inches(3.15))
        frame.fill.background()
        frame.line.color.rgb = RGBColor(*theme["secondary"])
        frame.line.width = Pt(3)

        print(f"  ✅ Картинка ({style})")
    except Exception as e:
        print(f"  ❌ {e}")

def create_conclusion_slide(prs, topic: str, theme: dict):
    """Финальный слайд"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*theme["secondary"])

    bg_shape = slide.shapes.add_shape(1, Inches(1.5), Inches(2), Inches(7), Inches(3.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(*theme["primary"])
    bg_shape.line.fill.background()

    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*theme["secondary"])
    p.alignment = PP_ALIGN.CENTER

    text_box2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(1))
    text_frame2 = text_box2.text_frame
    p = text_frame2.paragraphs[0]
    p.text = topic
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*theme["text_secondary"])
    p.alignment = PP_ALIGN.CENTER

def create_presentation(topic: str, num_slides: int, template: str, language: str = "russian") -> Presentation:
    """Создает презентацию с красивыми локальными картинками"""
    print(f"\n🎨 Создаю презентацию: {topic} ({num_slides} слайдов)")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme = THEMES.get(template, THEMES["modern"])

    styles = ["abstract", "tech", "organic", "cards"]

    if language == "russian":
        slides = [
            ("Введение", "• Основные концепции\n• Исторический контекст\n• Текущее состояние"),
            ("Ключевые моменты", "• Первый важный аспект\n• Второй ключевой момент\n• Третий важный пункт"),
            ("Практическое применение", "• Использование в проектах\n• Примеры успеха\n• Практические советы"),
            ("Преимущества", "• Повышение эффективности\n• Снижение затрат\n• Масштабируемость"),
            ("Вызовы и решения", "• Технические трудности\n• Стратегии преодоления\n• Инновационные подходы"),
            ("Будущее развития", "• Новые направления\n• Прогнозы развития\n• Возможности роста"),
            ("Детальный анализ", "• Углубленное изучение\n• Статистика и данные\n• Сравнительный анализ"),
            ("Рекомендации", "• Лучшие практики\n• Стратегические шаги\n• Реализация плана"),
        ]
    else:
        slides = [
            ("Introduction", "• Main concepts\n• Historical context\n• Current state"),
            ("Key Points", "• First aspect\n• Second moment\n• Third point"),
            ("Practical Applications", "• Project usage\n• Success examples\n• Practical tips"),
            ("Benefits", "• Increased efficiency\n• Cost reduction\n• Scalability"),
            ("Challenges and Solutions", "• Technical difficulties\n• Strategies\n• Innovative approaches"),
            ("Future Development", "• New directions\n• Development forecast\n• Growth opportunities"),
            ("Detailed Analysis", "• In-depth study\n• Statistics and data\n• Comparative analysis"),
            ("Recommendations", "• Best practices\n• Strategic steps\n• Implementation plan"),
        ]

    create_title_slide(prs, topic, theme)
    print(f"✅ Слайд 1: Титул")

    slides_to_create = min(num_slides - 2, len(slides))
    for i in range(slides_to_create):
        title, content = slides[i]
        style = styles[i % len(styles)]
        print(f"\n📄 Слайд {i + 2}: {title}")
        create_content_slide(prs, title, content, seed=i+42, theme=theme, style=style)

    create_conclusion_slide(prs, topic, theme)
    print(f"\n✅ Слайд {num_slides}: Спасибо")

    return prs

@app.route('/', methods=['GET'])
def health_check():
    return jsonify({'status': 'ok', 'message': '✅ Backend работает!'})

@app.route('/api/generate', methods=['POST'])
def generate_presentation():
    print(f"\n{'='*60}")
    print(f"📨 Запрос в {datetime.now().strftime('%H:%M:%S')}")
    print("="*60)

    try:
        data = request.json or {}

        topic = data.get('topic', 'Presentation').strip()
        num_slides = int(data.get('slides', 7))
        template = data.get('template', 'modern')
        output_path = data.get('outputPath', './presentations').strip()
        language = data.get('language', 'russian')

        if not topic:
            return jsonify({'success': False, 'error': 'Тема не может быть пустой'}), 400

        num_slides = max(3, min(num_slides, 15))
        template = template if template in THEMES else 'modern'

        print(f"📋 Параметры: {topic}, {num_slides} слайдов, {template}, {language}")

        os.makedirs(output_path, exist_ok=True)

        prs = create_presentation(topic, num_slides, template, language)

        safe_topic = "".join(c if c.isalnum() or c in (' ', '-', '_') else '' for c in topic)
        filename = f"{safe_topic.replace(' ', '_')}.pptx"
        filepath = os.path.join(output_path, filename)

        try:
            prs.save(filepath)
        except PermissionError:
            print(f"⚠️  Файл открыт, использую временное имя...")
            filename = f"presentation_{int(time.time())}.pptx"
            filepath = os.path.join(output_path, filename)
            prs.save(filepath)

        print(f"\n✅ ГОТОВО!")
        print(f"📁 {filepath}")
        print("="*60)

        return jsonify({
            'success': True,
            'message': f'✅ Презентация "{filename}" готова!',
            'filepath': os.path.abspath(filepath),
            'filename': filename,
            'slides_count': num_slides
        })

    except Exception as e:
        print(f"\n❌ ОШИБКА: {e}")
        print(traceback.format_exc())
        print("="*60)
        return jsonify({'success': False, 'error': str(e)}), 500

if __name__ == '__main__':
    print("\n🎨 ГЕНЕРАТОР ПРЕЗЕНТАЦИЙ v7.0")
    print("📍 http://127.0.0.1:5000")
    print("✨ Красивые слайды с локальными картинками!\n")
    app.run(debug=True, port=5000, host='127.0.0.1')
